from django.shortcuts import render
from django.shortcuts import redirect
from django.conf import settings
from django.core.files.storage import FileSystemStorage
from django.views.decorators.http import require_POST
import requests
import time
from azure.core.credentials import AzureKeyCredential
from azure.ai.textanalytics import TextAnalyticsClient, ExtractSummaryAction
import math
import numpy as np
import tensorflow as tf
from tensorflow import keras
import keras_cv
import requests
import urllib.request
from PIL import Image
from io import BytesIO
import marshal
import types
import random
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import pptx
import math


def index(request):
    context = {}

    if request.method == "POST" and request.FILES["myfile"]:
        myfile = request.FILES["myfile"]
        fs = FileSystemStorage()
        filename = fs.save(myfile.name, myfile)
        uploaded_file_url = fs.url(filename)
        context.update({"uploaded_file_url": uploaded_file_url})

    return render(request, "index.html", context=context)


@require_POST
def receive_audio(request):
    chapters, auto_highlights_result = audio_summarization(
        request.POST.get("audio_link")
    )
    keywords = []
    for x in auto_highlights_result["results"]:
        keywords.append(x["text"])
    summary = [x for x in map(str.strip, chapters.split("- ")) if x]

    keywords_mapping = {}
    author = request.POST.get("author")
    heading = request.POST.get("title")

    for keyword in keywords:
        if next((s for s in summary if keyword in s), None) != None:
            keywords_mapping[keyword] = summary.index(
                next((s for s in summary if keyword in s), None)
            )

    if len(keywords_mapping) > 5:
        keywords_mapping = dict(random.sample(list(keywords_mapping.items()), 5))

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = heading
    subtitle.text = author

    for i in range(math.ceil(len(summary) / 3)):
        j = i + 1

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Summary"
        tf = body_shape.text_frame
        print(summary[j * 3 - 3])
        if len(summary) > j * 3 - 1:
            tf.text = (
                summary[j * 3 - 3]
                + "\n"
                + summary[j * 3 - 2]
                + "\n"
                + summary[j * 3 - 1]
            )
        elif len(summary) > j * 3 - 2:
            tf.text = summary[j * 3 - 3] + "\n" + summary[j * 3 - 2]
        else:
            tf.text = summary[j * 3 - 3]
        for paragraph in body_shape.text_frame.paragraphs:
            paragraph.font.size = pptx.util.Pt(22)

        for x in list(keywords_mapping.values()):
            if j * 3 - 3 <= x <= j * 3 - 1:
                particular_keyword = list(keywords_mapping.keys())[
                    list(keywords_mapping.values()).index(x)
                ]
                pic_left = int(prs.slide_width * 0.15)
                pic_top = int(prs.slide_height * 0.1)
                pic_width = int(prs.slide_width * 0.7)
                pic_height = int(pic_width * 512 / 512)
                slide = prs.slides.add_slide(blank_slide_layout)
                tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
                p = tb.text_frame.add_paragraph()
                p.text = particular_keyword
                p.font.size = pptx.util.Pt(22)
                left = top = Inches(1.75)
                try:
                    image = generate_image(particular_keyword)
                    pic = slide.shapes.add_picture(image, left, top, height=Inches(5))
                except:
                    pass

    prs.save("result.pptx")

    return redirect("generated")


@require_POST
def receive_text(request):
    try:
        myfile = request.FILES["myfile"]
        with open(myfile.name, "r") as f:
            summary, keywords = sample_extractive_summarization([f.read().rstrip()])
    except:
        pass

    keywords_mapping = {}
    author = request.POST.get("author")
    heading = request.POST.get("title")
    summary = [x for x in map(str.strip, summary.split(".")) if x]
    keywords = [x for x in map(str.strip, keywords.split(",")) if x]

    print(summary)
    print(keywords)

    for keyword in keywords:
        if next((s for s in summary if keyword in s), None) != None:
            keywords_mapping[keyword] = summary.index(
                next((s for s in summary if keyword in s), None)
            )

    if len(keywords_mapping) > 5:
        keywords_mapping = dict(random.sample(list(keywords_mapping.items()), 5))

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = heading
    subtitle.text = author

    for i in range(math.ceil(len(summary) / 3) - 2):
        j = i + 1

        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = "Summary"
        tf = body_shape.text_frame
        tf.text = (
            summary[j * 3 - 3] + "\n" + summary[j * 3 - 2] + "\n" + summary[j * 3 - 1]
        )
        for paragraph in body_shape.text_frame.paragraphs:
            paragraph.font.size = pptx.util.Pt(22)

        for x in list(keywords_mapping.values()):
            if j * 3 - 3 <= x <= j * 3 - 1:
                particular_keyword = list(keywords_mapping.keys())[
                    list(keywords_mapping.values()).index(x)
                ]
                pic_left = int(prs.slide_width * 0.15)
                pic_top = int(prs.slide_height * 0.1)
                pic_width = int(prs.slide_width * 0.7)
                pic_height = int(pic_width * 512 / 512)
                slide = prs.slides.add_slide(blank_slide_layout)
                tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
                p = tb.text_frame.add_paragraph()
                p.text = particular_keyword
                p.font.size = pptx.util.Pt(22)
                left = top = Inches(1.75)
                try:
                    image = generate_image(particular_keyword)
                    pic = slide.shapes.add_picture(image, left, top, height=Inches(5))
                except:
                    pass

    prs.save("result.pptx")

    return redirect("generated")


def sample_extractive_summarization(document):
    """Get summarization and keywords of .txt file"""
    endpoint = "https://hack121022.cognitiveservices.azure.com/"
    key = settings.AZURE_API_KEY

    text_analytics_client = TextAnalyticsClient(
        endpoint=endpoint,
        credential=AzureKeyCredential(key),
    )

    poller = text_analytics_client.begin_analyze_actions(
        document,
        actions=[
            ExtractSummaryAction(
                max_sentence_count=min(20, int(len(document[0]) * 0.1))
            ),
        ],
    )

    summary, keywords = "", ""

    document_results = poller.result()
    for extract_summary_results in document_results:
        for result in extract_summary_results:
            # if result.kind == "ExtractiveSummarization":
            print(
                "Summary extracted: \n{}".format(
                    " ".join([sentence.text for sentence in result.sentences])
                )
            )
            summary = " ".join([sentence.text for sentence in result.sentences])
        # elif result.is_error is True:
        #     print("...Is an error with code '{}' and message '{}'".format(
        #         result.code, result.message
        #     ))

    result = text_analytics_client.extract_key_phrases(document)
    for idx, doc in enumerate(result):
        if not doc.is_error:
            print(
                "Key phrases in article #{}: {}".format(
                    idx + 1, ", ".join(doc.key_phrases)
                )
            )
            keywords = ", ".join(doc.key_phrases)

    return summary, keywords


def audio_summarization(audio_link):
    """Get summarization of audio using AssemblyAI API"""
    import pprint
    import requests
    import time

    print(audio_link)
    endpoint = "https://api.assemblyai.com/v2/transcript"
    json = {
        "audio_url": audio_link,
        "auto_highlights": True,
        "summarization": True,
        "summary_model": "informative",
        "summary_type": "bullets",
    }
    headers = {
        "authorization": settings.ASSEMBLYAI_API_KEY,
        "content-type": "application/json",
    }
    response = requests.post(endpoint, json=json, headers=headers)

    endpoint = "https://api.assemblyai.com/v2/transcript/" + response.json()["id"]
    headers = {
        "authorization": settings.ASSEMBLYAI_API_KEY,
    }
    response = requests.get(endpoint, headers=headers)

    while response.json()["status"] != "completed":
        time.sleep(2)
        response = requests.get(endpoint, headers=headers)

    pprint.pprint(response.json()["summary"])
    print()
    pprint.pprint(response.json()["auto_highlights_result"])
    print()
    pprint.pprint(response.json()["status"])

    return response.json()["summary"], response.json()["auto_highlights_result"]


def generate_image(keyword):
    with open("serialized_bin", "rb") as f:
        serialized = marshal.loads(f.read())
        predict = types.FunctionType(serialized, globals(), "predict")

    myfile = predict(keyword)
    name = "{}.jpg".format(keyword)
    myfile.save(name)

    return name
