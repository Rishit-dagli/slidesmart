import random
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import pptx
import math

keywords = ["abc", "def"]
keywords_mapping = {}
author = "Author"
heading = "Title"
summary = ["fabc", "rksbv", "diuf", "fdef", "erjhf", "rf"]
# summary = [x for x in map(str.strip, summary.split('.')) if x]

for keyword in keywords:
    keywords_mapping[keyword] = summary.index(
        next((s for s in summary if keyword in s), None)
    )

if len(keywords_mapping) > 5:
    keywords_mapping = dict(random.sample(list(keywords_mapping.items()), 5))

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]

slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = heading
subtitle.text = author

for i in range(math.ceil(len(summary) / 3)):
    j = i + 1

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Adding a Bullet Slide"
    tf = body_shape.text_frame
    tf.text = summary[j * 3 - 3] + "\n" + summary[j * 3 - 2] + "\n" + summary[j * 3 - 1]

    for x in list(keywords_mapping.values()):
        if j * 3 - 3 <= x <= j * 3 - 1:
            particular_keyword = list(keywords_mapping.keys())[
                list(keywords_mapping.values()).index(x)
            ]
            pic_left = int(prs.slide_width * 0.15)
            pic_top = int(prs.slide_height * 0.1)
            pic_width = int(prs.slide_width * 0.7)
            pic_height = int(pic_width * 512 / 512)
            slide = prs.slides.add_slide(blank_slide_layout)
            tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
            p = tb.text_frame.add_paragraph()
            p.text = particular_keyword
            p.font.size = pptx.util.Pt(22)
            left = top = Inches(1.75)
            pic = slide.shapes.add_picture("abc.png", left, top, height=Inches(5))

prs.save("test.pptx")
